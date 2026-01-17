# ADIS Document Parser Azure Function
# Version: 1.0
# Purpose: Parse uploaded documents (Excel, Word, PDF, PowerPoint) into structured data

import azure.functions as func
import logging
import json
import hashlib
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple
import io

# Document parsing libraries
import pandas as pd
from openpyxl import load_workbook
import docx
import pptx
try:
    import fitz  # PyMuPDF for PDF parsing
except ImportError:
    fitz = None

app = func.FunctionApp()

# =============================================================================
# CONFIGURATION
# =============================================================================

MAX_FILE_SIZE_MB = 250
MAX_ROWS = 1_000_000
SAMPLE_SIZE = 10  # Number of sample values to store per column

# Semantic type detection patterns
SEMANTIC_PATTERNS = {
    'CUSTOMER_ID': {
        'column_patterns': ['customer_id', 'cust_id', 'client_id', 'member_id', 'user_id', 'account_id'],
        'value_patterns': [r'^[A-Z]{2,4}[-_]?\d{4,10}$', r'^CUST[-_]\d+$'],
        'characteristics': {'unique_ratio': 0.9, 'non_null_ratio': 0.95}
    },
    'TRANSACTION_ID': {
        'column_patterns': ['transaction_id', 'txn_id', 'order_id', 'invoice_id', 'receipt_id'],
        'value_patterns': [r'^TXN[-_]?\d+$', r'^ORD[-_]?\d+$'],
        'characteristics': {'unique_ratio': 0.99}
    },
    'EMAIL': {
        'column_patterns': ['email', 'email_address', 'e_mail'],
        'value_patterns': [r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'],
        'characteristics': {}
    },
    'PHONE': {
        'column_patterns': ['phone', 'telephone', 'mobile', 'cell', 'contact_number'],
        'value_patterns': [r'^\+?1?[-.\s]?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}$'],
        'characteristics': {}
    },
    'DATE': {
        'column_patterns': ['date', 'created', 'updated', 'timestamp', 'time', 'datetime'],
        'value_patterns': [],  # Use pandas date parsing
        'characteristics': {}
    },
    'CURRENCY': {
        'column_patterns': ['amount', 'price', 'cost', 'revenue', 'spend', 'total', 'value', 'ltv', 'monetary'],
        'value_patterns': [r'^\$?[\d,]+\.?\d*$'],
        'characteristics': {'dtype': 'numeric'}
    },
    'CATEGORY': {
        'column_patterns': ['category', 'type', 'status', 'tier', 'segment', 'group', 'class'],
        'value_patterns': [],
        'characteristics': {'unique_ratio_max': 0.1, 'dtype': 'string'}
    },
    'ENGAGEMENT_SCORE': {
        'column_patterns': ['score', 'rating', 'engagement', 'loyalty_score', 'nps'],
        'value_patterns': [],
        'characteristics': {'dtype': 'numeric', 'min_val': 0, 'max_val': 100}
    },
    'LOYALTY_TIER': {
        'column_patterns': ['tier', 'level', 'membership', 'loyalty_tier', 'status'],
        'value_patterns': [r'^(gold|silver|bronze|platinum|diamond)$'],
        'characteristics': {'unique_ratio_max': 0.01}
    }
}

# =============================================================================
# MAIN FUNCTION ENDPOINT
# =============================================================================

@app.route(route="parse-document", methods=["POST"])
async def parse_document(req: func.HttpRequest) -> func.HttpResponse:
    """
    Main endpoint for document parsing.
    Accepts file upload and returns parsed data with inferred schema.
    """
    logging.info('ADIS Document Parser: Processing request')
    
    try:
        # Get file from request
        file_content = req.get_body()
        file_name = req.headers.get('X-File-Name', 'unknown')
        upload_job_id = req.headers.get('X-Upload-Job-Id', str(hashlib.md5(file_content[:1000]).hexdigest()))
        
        if not file_content:
            return func.HttpResponse(
                json.dumps({'error': 'No file content provided'}),
                status_code=400,
                mimetype='application/json'
            )
        
        # Validate file size
        file_size_mb = len(file_content) / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            return func.HttpResponse(
                json.dumps({'error': f'File size {file_size_mb:.1f}MB exceeds maximum {MAX_FILE_SIZE_MB}MB'}),
                status_code=400,
                mimetype='application/json'
            )
        
        # Determine file type and parse
        file_extension = file_name.lower().split('.')[-1] if '.' in file_name else ''
        
        parser_map = {
            'xlsx': parse_excel,
            'xls': parse_excel,
            'csv': parse_csv,
            'docx': parse_word,
            'pdf': parse_pdf,
            'pptx': parse_powerpoint
        }
        
        parser = parser_map.get(file_extension)
        if not parser:
            return func.HttpResponse(
                json.dumps({'error': f'Unsupported file type: {file_extension}'}),
                status_code=400,
                mimetype='application/json'
            )
        
        # Parse the document
        parse_result = parser(file_content, file_name)
        
        if parse_result.get('error'):
            return func.HttpResponse(
                json.dumps(parse_result),
                status_code=400,
                mimetype='application/json'
            )
        
        # Infer schema for each dataframe
        schemas = []
        for df_info in parse_result.get('dataframes', []):
            df = df_info['dataframe']
            schema = infer_schema(df, df_info.get('sheet_name', 'default'))
            schemas.append(schema)
        
        # Build response
        response = {
            'upload_job_id': upload_job_id,
            'file_name': file_name,
            'file_type': file_extension.upper(),
            'file_size_bytes': len(file_content),
            'parse_timestamp': datetime.utcnow().isoformat(),
            'status': 'SUCCESS',
            'summary': {
                'total_sheets': len(parse_result.get('dataframes', [])),
                'total_rows': sum(df['row_count'] for df in parse_result.get('dataframes', [])),
                'total_columns': sum(len(s['columns']) for s in schemas)
            },
            'schemas': schemas,
            'data_preview': [
                {
                    'sheet_name': df_info.get('sheet_name', 'Sheet1'),
                    'preview': df_info['dataframe'].head(5).to_dict(orient='records')
                }
                for df_info in parse_result.get('dataframes', [])[:3]  # Limit preview to first 3 sheets
            ]
        }
        
        return func.HttpResponse(
            json.dumps(response, default=str),
            status_code=200,
            mimetype='application/json'
        )
        
    except Exception as e:
        logging.error(f'ADIS Document Parser error: {str(e)}')
        return func.HttpResponse(
            json.dumps({'error': str(e), 'status': 'FAILED'}),
            status_code=500,
            mimetype='application/json'
        )


# =============================================================================
# FILE PARSERS
# =============================================================================

def parse_excel(content: bytes, file_name: str) -> Dict[str, Any]:
    """Parse Excel file (xlsx, xls) into dataframes."""
    try:
        excel_file = io.BytesIO(content)
        
        # Read all sheets
        xlsx = pd.ExcelFile(excel_file)
        dataframes = []
        
        for sheet_name in xlsx.sheet_names:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            
            if len(df) > MAX_ROWS:
                logging.warning(f'Sheet {sheet_name} has {len(df)} rows, truncating to {MAX_ROWS}')
                df = df.head(MAX_ROWS)
            
            dataframes.append({
                'sheet_name': sheet_name,
                'dataframe': df,
                'row_count': len(df),
                'column_count': len(df.columns)
            })
        
        return {'dataframes': dataframes}
        
    except Exception as e:
        return {'error': f'Excel parsing failed: {str(e)}'}


def parse_csv(content: bytes, file_name: str) -> Dict[str, Any]:
    """Parse CSV file into dataframe."""
    try:
        # Try different encodings
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                df = pd.read_csv(io.BytesIO(content), encoding=encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            return {'error': 'Unable to decode CSV with supported encodings'}
        
        if len(df) > MAX_ROWS:
            logging.warning(f'CSV has {len(df)} rows, truncating to {MAX_ROWS}')
            df = df.head(MAX_ROWS)
        
        return {
            'dataframes': [{
                'sheet_name': 'Sheet1',
                'dataframe': df,
                'row_count': len(df),
                'column_count': len(df.columns)
            }]
        }
        
    except Exception as e:
        return {'error': f'CSV parsing failed: {str(e)}'}


def parse_word(content: bytes, file_name: str) -> Dict[str, Any]:
    """Parse Word document, extracting tables."""
    try:
        doc = docx.Document(io.BytesIO(content))
        dataframes = []
        
        # Extract tables from Word document
        for i, table in enumerate(doc.tables):
            data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                data.append(row_data)
            
            if len(data) > 1:  # At least header + 1 data row
                df = pd.DataFrame(data[1:], columns=data[0])
                dataframes.append({
                    'sheet_name': f'Table_{i+1}',
                    'dataframe': df,
                    'row_count': len(df),
                    'column_count': len(df.columns)
                })
        
        if not dataframes:
            return {'error': 'No tables found in Word document'}
        
        return {'dataframes': dataframes}
        
    except Exception as e:
        return {'error': f'Word parsing failed: {str(e)}'}


def parse_pdf(content: bytes, file_name: str) -> Dict[str, Any]:
    """Parse PDF document, extracting tables."""
    if fitz is None:
        return {'error': 'PDF parsing not available (PyMuPDF not installed)'}
    
    try:
        pdf = fitz.open(stream=content, filetype='pdf')
        dataframes = []
        
        for page_num, page in enumerate(pdf):
            # Extract tables using PyMuPDF's table detection
            tables = page.find_tables()
            
            for i, table in enumerate(tables):
                df = table.to_pandas()
                if len(df) > 0:
                    dataframes.append({
                        'sheet_name': f'Page{page_num+1}_Table{i+1}',
                        'dataframe': df,
                        'row_count': len(df),
                        'column_count': len(df.columns)
                    })
        
        if not dataframes:
            return {'error': 'No tables found in PDF document'}
        
        return {'dataframes': dataframes}
        
    except Exception as e:
        return {'error': f'PDF parsing failed: {str(e)}'}


def parse_powerpoint(content: bytes, file_name: str) -> Dict[str, Any]:
    """Parse PowerPoint document, extracting tables."""
    try:
        prs = pptx.Presentation(io.BytesIO(content))
        dataframes = []
        
        for slide_num, slide in enumerate(prs.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if shape.has_table:
                    table = shape.table
                    data = []
                    
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        data.append(row_data)
                    
                    if len(data) > 1:
                        df = pd.DataFrame(data[1:], columns=data[0])
                        dataframes.append({
                            'sheet_name': f'Slide{slide_num+1}_Table{shape_num+1}',
                            'dataframe': df,
                            'row_count': len(df),
                            'column_count': len(df.columns)
                        })
        
        if not dataframes:
            return {'error': 'No tables found in PowerPoint document'}
        
        return {'dataframes': dataframes}
        
    except Exception as e:
        return {'error': f'PowerPoint parsing failed: {str(e)}'}


# =============================================================================
# SCHEMA INFERENCE
# =============================================================================

def infer_schema(df: pd.DataFrame, sheet_name: str) -> Dict[str, Any]:
    """
    Infer semantic schema for a dataframe.
    Returns column-level metadata with detected types and characteristics.
    """
    schema = {
        'sheet_name': sheet_name,
        'row_count': len(df),
        'column_count': len(df.columns),
        'columns': []
    }
    
    for col_idx, col_name in enumerate(df.columns):
        col_data = df[col_name]
        
        column_schema = {
            'column_name': str(col_name),
            'column_index': col_idx,
            'detected_type': detect_semantic_type(col_name, col_data),
            'confidence_score': 0.0,  # Will be set by detect_semantic_type
            'pandas_dtype': str(col_data.dtype),
            'null_percentage': round((col_data.isna().sum() / len(col_data)) * 100, 2),
            'unique_count': col_data.nunique(),
            'unique_ratio': round(col_data.nunique() / len(col_data), 4) if len(col_data) > 0 else 0,
            'sample_values': get_sample_values(col_data),
            'statistics': get_column_statistics(col_data)
        }
        
        # Set confidence based on detection method
        column_schema['confidence_score'] = calculate_confidence(column_schema)
        
        # Determine RFM role if applicable
        column_schema['rfm_role'] = detect_rfm_role(col_name, column_schema['detected_type'], col_data)
        
        schema['columns'].append(column_schema)
    
    # Add data quality summary
    schema['data_quality'] = {
        'completeness': round(100 - df.isna().sum().sum() / (len(df) * len(df.columns)) * 100, 2),
        'potential_duplicates': df.duplicated().sum(),
        'has_customer_key': any(c['detected_type'] == 'CUSTOMER_ID' for c in schema['columns']),
        'has_transaction_data': any(c['detected_type'] in ['TRANSACTION_ID', 'DATE', 'CURRENCY'] for c in schema['columns']),
        'rfm_ready': is_rfm_ready(schema['columns'])
    }
    
    return schema


def detect_semantic_type(col_name: str, col_data: pd.Series) -> str:
    """Detect the semantic type of a column based on name and values."""
    col_name_lower = col_name.lower().replace(' ', '_').replace('-', '_')
    
    # Check each semantic type
    for sem_type, patterns in SEMANTIC_PATTERNS.items():
        # Check column name patterns
        if any(pattern in col_name_lower for pattern in patterns['column_patterns']):
            return sem_type
        
        # Check value patterns (sample first 100 non-null values)
        if patterns.get('value_patterns'):
            import re
            sample = col_data.dropna().head(100).astype(str)
            for pattern in patterns['value_patterns']:
                matches = sample.str.match(pattern, case=False).sum()
                if matches / len(sample) > 0.8:  # 80% match threshold
                    return sem_type
    
    # Check for date type using pandas
    if col_data.dtype == 'datetime64[ns]':
        return 'DATE'
    
    try:
        pd.to_datetime(col_data.dropna().head(100), errors='raise')
        return 'DATE'
    except:
        pass
    
    # Check for numeric
    if pd.api.types.is_numeric_dtype(col_data):
        return 'NUMERIC'
    
    # Default to TEXT
    return 'TEXT'


def detect_rfm_role(col_name: str, detected_type: str, col_data: pd.Series) -> str:
    """Detect if column should be used for RFM analysis."""
    col_name_lower = col_name.lower()
    
    if detected_type == 'CUSTOMER_ID':
        return 'Customer Key'
    
    if detected_type == 'DATE':
        if any(kw in col_name_lower for kw in ['transaction', 'purchase', 'order', 'created']):
            return 'Recency Date'
    
    if detected_type == 'CURRENCY':
        if any(kw in col_name_lower for kw in ['amount', 'total', 'revenue', 'value', 'spend']):
            return 'Monetary Value'
    
    # Frequency is typically derived, not a direct column
    
    return 'None'


def is_rfm_ready(columns: List[Dict]) -> bool:
    """Check if dataset has minimum columns for RFM analysis."""
    roles = [c.get('rfm_role', 'None') for c in columns]
    return 'Customer Key' in roles and 'Recency Date' in roles


def get_sample_values(col_data: pd.Series) -> List[str]:
    """Get sample non-null values for preview."""
    samples = col_data.dropna().head(SAMPLE_SIZE).tolist()
    return [str(s)[:100] for s in samples]  # Truncate long values


def get_column_statistics(col_data: pd.Series) -> Dict[str, Any]:
    """Calculate basic statistics for a column."""
    stats = {}
    
    if pd.api.types.is_numeric_dtype(col_data):
        stats['min'] = float(col_data.min()) if not col_data.isna().all() else None
        stats['max'] = float(col_data.max()) if not col_data.isna().all() else None
        stats['mean'] = float(col_data.mean()) if not col_data.isna().all() else None
        stats['median'] = float(col_data.median()) if not col_data.isna().all() else None
        stats['std'] = float(col_data.std()) if not col_data.isna().all() else None
    elif col_data.dtype == 'datetime64[ns]':
        stats['min'] = str(col_data.min()) if not col_data.isna().all() else None
        stats['max'] = str(col_data.max()) if not col_data.isna().all() else None
    else:
        # String column - get value distribution
        value_counts = col_data.value_counts().head(10)
        stats['top_values'] = value_counts.to_dict()
    
    return stats


def calculate_confidence(column_schema: Dict) -> float:
    """Calculate confidence score for type detection."""
    base_confidence = 0.6
    
    detected_type = column_schema['detected_type']
    
    # Boost for low null percentage
    if column_schema['null_percentage'] < 5:
        base_confidence += 0.1
    
    # Boost for column name match
    col_name_lower = column_schema['column_name'].lower()
    if detected_type in SEMANTIC_PATTERNS:
        patterns = SEMANTIC_PATTERNS[detected_type]['column_patterns']
        if any(p in col_name_lower for p in patterns):
            base_confidence += 0.2
    
    # Boost for appropriate unique ratio
    if detected_type in ['CUSTOMER_ID', 'TRANSACTION_ID']:
        if column_schema['unique_ratio'] > 0.9:
            base_confidence += 0.1
    elif detected_type == 'CATEGORY':
        if column_schema['unique_ratio'] < 0.1:
            base_confidence += 0.1
    
    return min(base_confidence, 1.0)


# =============================================================================
# HEALTH CHECK ENDPOINT
# =============================================================================

@app.route(route="health", methods=["GET"])
def health_check(req: func.HttpRequest) -> func.HttpResponse:
    """Health check endpoint."""
    return func.HttpResponse(
        json.dumps({
            'status': 'healthy',
            'service': 'ADIS Document Parser',
            'version': '1.0',
            'timestamp': datetime.utcnow().isoformat(),
            'capabilities': ['xlsx', 'xls', 'csv', 'docx', 'pdf', 'pptx']
        }),
        status_code=200,
        mimetype='application/json'
    )
